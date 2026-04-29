# Create a PowerPoint using python-pptx
from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

def add_slide(title, bullets):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = b
        p.level = 0

slides_content = [
("Explainable AI (XAI) - Artigo Científico",
 ["A Multidisciplinary Survey and Framework for Design and Evaluation of Explainable AI Systems",
  "Aluno: José Borges",
  "Disciplina: IHC"]),

("Introdução",
 ["Crescimento da Inteligência Artificial",
  "Uso em decisões importantes",
  "Problema: falta de transparência"]),

("Problema",
 ["Sistemas “caixa preta”",
  "Falta de explicação",
  "Risco de decisões injustas"]),

("O que é XAI",
 ["Explainable AI (XAI)",
  "Sistemas que explicam decisões",
  "Foco no usuário"]),

("Objetivo do artigo",
 ["Revisar pesquisas existentes",
  "Organizar métodos",
  "Criar um framework"]),

("Abordagem",
 ["Análise de mais de 200 artigos",
  "Diferentes áreas",
  "Comparação de métodos"]),

("Áreas envolvidas",
 ["Machine Learning",
  "IHC",
  "Visualização",
  "Psicologia"]),

("Tipos de explicação",
 ["Why (por quê)",
  "How (como funciona)",
  "What-if (e se...)"]),

("Tipos de usuários",
 ["Usuários comuns",
  "Especialistas em dados",
  "Especialistas em IA"]),

("Objetivos do XAI",
 ["Transparência",
  "Confiança",
  "Redução de viés",
  "Compreensão"]),

("Avaliação",
 ["Entendimento do usuário",
  "Confiança",
  "Desempenho",
  "Satisfação"]),

("Framework",
 ["Guia de design",
  "Integra diferentes áreas",
  "Apoia desenvolvimento de XAI"]),

("Contribuições",
 ["Organização do conhecimento",
  "Integração de áreas",
  "Guia prático"]),

("Limitações",
 ["Área ainda em desenvolvimento",
  "Difícil padronizar explicações"]),

("Conclusão",
 ["XAI é essencial",
  "Impacto direto na sociedade",
  "Futuro da IA"])
]

for title, bullets in slides_content:
    add_slide(title, bullets)

file_path = "/mnt/data/slides_XAI.pptx"
prs.save(file_path)

file_path